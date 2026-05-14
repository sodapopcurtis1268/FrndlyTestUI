"""
Builds the Playwright automation presentation as a .pptx file.
Run: python3 build_deck.py
Output: FrndlyTV_Automation_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
FRNDLY_BLUE   = RGBColor(0x00, 0x6F, 0xC6)   # Frndly TV primary blue
FRNDLY_DARK   = RGBColor(0x0A, 0x1A, 0x3A)   # Deep navy (backgrounds)
FRNDLY_GREEN  = RGBColor(0x00, 0xC4, 0x7F)   # Accent green (success / go)
FRNDLY_YELLOW = RGBColor(0xFF, 0xC8, 0x00)   # Accent yellow (highlight)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY      = RGBColor(0x6B, 0x7B, 0x8D)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank — full design control


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()   # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bullet_box(slide, items, left, top, width, height,
               size=20, color=WHITE, dot_color=None, spacing=None):
    """Add a bulleted list with custom dot colour."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if spacing:
            p.space_before = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT

        # coloured bullet dot
        if dot_color:
            dot = p.add_run()
            dot.text = "● "
            dot.font.size = Pt(size - 2)
            dot.font.color.rgb = dot_color
            dot.font.bold = True

        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color


def accent_bar(slide, color=FRNDLY_GREEN):
    """Thin horizontal accent line near the top."""
    rect(slide, 0, 0.95, 13.33, 0.06, fill=color)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)          # full background
rect(s, 0, 0, 13.33, 0.12, fill=FRNDLY_BLUE)          # top stripe
rect(s, 0, 7.38, 13.33, 0.12, fill=FRNDLY_BLUE)       # bottom stripe
rect(s, 0.4, 2.8, 12.5, 0.08, fill=FRNDLY_GREEN)      # green divider

txt(s, "FRNDLY TV", 0.5, 1.1, 12, 1,
    size=22, bold=True, color=FRNDLY_GREEN, align=PP_ALIGN.CENTER)
txt(s, "Test Automation", 0.5, 1.9, 12, 1.2,
    size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Catching bugs before your customers do", 0.5, 3.1, 12, 0.8,
    size=26, color=FRNDLY_YELLOW, align=PP_ALIGN.CENTER)
txt(s, "Powered by Playwright  ·  GitHub Actions  ·  TypeScript",
    0.5, 6.5, 12, 0.6, size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem (why we're here)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_YELLOW)

txt(s, "The Problem We're Solving", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

problems = [
    "A customer finds a broken feature before we do",
    "Manual testing can't keep up with every release",
    "The same checks get repeated by hand — every time",
    "One missed test = one support ticket, one frustrated subscriber",
    "We have no way to know if a new update broke something at 2 AM",
]
bullet_box(s, problems, 0.8, 2.1, 11.5, 4.5,
           size=22, color=WHITE, dot_color=FRNDLY_YELLOW, spacing=6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Goal
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_GREEN)

txt(s, "What We Want to Achieve", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

# 3 goal cards
card_data = [
    ("🛡️", "Catch Bugs First", "Find broken features before\ncustomers see them"),
    ("⚡", "Ship Faster", "Every code change tested\nautomatically in minutes"),
    ("📊", "Measure Quality", "Real numbers: how fast does\nvideo load? Is it trending worse?"),
]
for i, (icon, title, body) in enumerate(card_data):
    x = 0.4 + i * 4.3
    rect(s, x, 2.1, 4.0, 3.8, fill=FRNDLY_BLUE)
    txt(s, icon,  x + 0.2, 2.3, 3.6, 0.7, size=32, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.2, 3.1, 3.6, 0.6, size=20, bold=True,
        color=FRNDLY_YELLOW, align=PP_ALIGN.CENTER)
    txt(s, body,  x + 0.2, 3.8, 3.6, 1.5, size=17, color=WHITE,
        align=PP_ALIGN.CENTER)

txt(s, "Bottom line: we find the problem — not the subscriber.",
    0.5, 6.2, 12, 0.6, size=20, bold=True,
    color=FRNDLY_GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is Playwright?
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_BLUE)

txt(s, "What is Playwright?", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

txt(s, "Think of it as a robot that operates your website\nexactly like a real customer — 24 hours a day, 7 days a week.",
    0.6, 2.0, 8.5, 1.4, size=22, color=FRNDLY_YELLOW)

facts = [
    "Built by Microsoft  —  the same company behind Windows & Office",
    "Used by Netflix, Spotify, GitHub, and thousands of other companies",
    "Records a VIDEO of every test so you can watch exactly what happened",
    "Completely FREE  —  no per-test charges, no cloud subscription",
    "Runs automatically every time a developer makes a change",
]
bullet_box(s, facts, 0.8, 3.6, 11.5, 3.5,
           size=20, color=WHITE, dot_color=FRNDLY_GREEN, spacing=4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What the Robot Does (visual flow)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s)

txt(s, "What the Robot Does — Step by Step", 0.5, 1.1, 12, 0.8,
    size=34, bold=True, color=WHITE)

steps = [
    ("1", "Opens the Frndly TV website"),
    ("2", "Logs in with a test account"),
    ("3", "Clicks a channel or movie"),
    ("4", "Measures exactly how many seconds until video plays"),
    ("5", "Takes a screenshot of the playing video"),
    ("6", "Navigates through Settings, signs out"),
    ("7", "Reports PASS ✅  or  FAIL ❌  with video evidence"),
]
for i, (num, step) in enumerate(steps):
    y = 2.1 + i * 0.67
    rect(s, 0.5, y, 0.55, 0.52, fill=FRNDLY_GREEN)
    txt(s, num, 0.5, y - 0.02, 0.55, 0.56,
        size=18, bold=True, color=FRNDLY_DARK, align=PP_ALIGN.CENTER)
    txt(s, step, 1.2, y, 11, 0.55, size=20, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What We've Built (the framework)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_YELLOW)

txt(s, "What We've Built for Frndly TV", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

# Left column
txt(s, "Smoke Tests  —  Run in ~2 minutes",
    0.5, 2.1, 6.0, 0.5, size=20, bold=True, color=FRNDLY_GREEN)
bullet_box(s, [
    "Login flow works end-to-end",
    "Live TV plays on the home page",
    "Random content row loads & streams",
    "Settings & Sign Out function correctly",
], 0.7, 2.7, 5.8, 2.2, size=18, color=WHITE, dot_color=FRNDLY_YELLOW)

# Right column
txt(s, "Regression Tests  —  One per content row",
    6.8, 2.1, 6.0, 0.5, size=20, bold=True, color=FRNDLY_BLUE)
bullet_box(s, [
    "17 individual row tests",
    "Fan Favorites, Staff Picks, Most Watched…",
    "Measures load time (TTFF) for each row",
    "Video recorded for every test run",
], 7.0, 2.7, 6.0, 2.2, size=18, color=WHITE, dot_color=FRNDLY_YELLOW)

rect(s, 6.5, 2.0, 0.05, 4.2, fill=MID_GRAY)  # divider

txt(s, "All tests run automatically on GitHub — FREE, no server needed.",
    0.5, 6.2, 12, 0.6, size=20, bold=True,
    color=FRNDLY_GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TTFF: The Business Metric
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_GREEN)

txt(s, "The Number That Matters Most", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

# Big TTFF callout
rect(s, 0.5, 2.1, 5.5, 3.8, fill=FRNDLY_BLUE)
txt(s, "TTFF", 0.5, 2.3, 5.5, 0.8,
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Time To First Frame", 0.5, 3.1, 5.5, 0.5,
    size=18, color=FRNDLY_YELLOW, align=PP_ALIGN.CENTER)
txt(s, "How many seconds from\n'I clicked play'\nto 'video is on my screen'",
    0.5, 3.7, 5.5, 1.5, size=20, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "We measure this every single test run.",
    0.5, 5.4, 5.5, 0.5, size=16, italic=True,
    color=FRNDLY_YELLOW, align=PP_ALIGN.CENTER)

# Right: why it matters
txt(s, "Why does this matter?", 6.3, 2.1, 6.7, 0.5,
    size=22, bold=True, color=FRNDLY_GREEN)
bullet_box(s, [
    "Subscribers cancel if video is slow to start",
    "We get an alert if load times get worse after a release",
    "We have proof — with a number — not just 'it felt slow'",
    "Typical result: Live TV starts in under 3 seconds",
], 6.3, 2.8, 6.7, 3.0, size=19, color=WHITE, dot_color=FRNDLY_YELLOW, spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GitHub Actions (the CI/CD angle)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_BLUE)

txt(s, "How Tests Run Automatically", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

# Flow diagram: 3 boxes + arrows
boxes = [
    ("Developer\nmakes a change", FRNDLY_BLUE),
    ("GitHub automatically\nruns all tests", FRNDLY_GREEN),
    ("Team gets a PASS\nor FAIL report", FRNDLY_YELLOW),
]
for i, (label, color) in enumerate(boxes):
    x = 0.4 + i * 4.3
    rect(s, x, 2.3, 3.9, 1.8, fill=color)
    txt(s, label, x, 2.3, 3.9, 1.8,
        size=20, bold=True, color=FRNDLY_DARK, align=PP_ALIGN.CENTER)
    if i < 2:
        txt(s, "→", x + 4.0, 2.8, 0.5, 0.8,
            size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_box(s, [
    "No one has to remember to run tests — it's automatic",
    "Tests complete in about 5 minutes",
    "Full video recordings saved for 30 days",
    "Works while the team sleeps — catches overnight issues",
    "Completely free using GitHub's servers",
], 0.8, 4.4, 11.5, 2.8,
   size=20, color=WHITE, dot_color=FRNDLY_GREEN, spacing=4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE DEMO placeholder
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
rect(s, 0, 0, 13.33, 0.12, fill=FRNDLY_GREEN)
rect(s, 0, 7.38, 13.33, 0.12, fill=FRNDLY_GREEN)

txt(s, "🎬  LIVE DEMO", 0.5, 1.8, 12, 1.2,
    size=54, bold=True, color=FRNDLY_GREEN, align=PP_ALIGN.CENTER)

demo_steps = [
    "1.  Open GitHub → Actions tab  (show the test history)",
    "2.  Trigger a manual run  →  watch tests execute in real time",
    "3.  Open the HTML report  →  show PASS / SKIP / FAIL per row",
    "4.  Click a test  →  play the recorded video of the robot browsing Frndly TV",
    "5.  Show the TTFF attachment  →  'Live Now loaded in 2.7 seconds'",
]
bullet_box(s, demo_steps, 1.0, 3.2, 11.3, 3.8,
           size=22, color=WHITE, dot_color=FRNDLY_YELLOW, spacing=6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DRM Explained Simply
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_YELLOW)

txt(s, "One Limitation We Hit — and How We Fix It", 0.5, 1.1, 12, 0.8,
    size=34, bold=True, color=WHITE)

# DRM explanation box
rect(s, 0.4, 2.0, 7.8, 2.3, fill=RGBColor(0x1A, 0x2A, 0x4A))
txt(s, "🔒  What is DRM?", 0.6, 2.1, 7.4, 0.6,
    size=22, bold=True, color=FRNDLY_YELLOW)
txt(s,
    "Movies and on-demand shows are protected by a digital lock called DRM.\n"
    "Your browser needs a special key from Frndly TV's server to unlock and play them.\n"
    "In an automated test robot, that key is sometimes refused — just like a bouncer\n"
    "checking IDs won't let someone in wearing a disguise.",
    0.6, 2.7, 7.5, 1.5, size=17, color=WHITE)

# What happens now
txt(s, "What happens today:", 0.5, 4.5, 6.0, 0.5,
    size=20, bold=True, color=FRNDLY_GREEN)
bullet_box(s, [
    "Live TV  →  ✅  Tests PASS  (no DRM on live streams)",
    "Movies / On-Demand  →  ⏭️  Tests SKIP with a clear note",
    "The robot says 'DRM blocked' — not 'something broke'",
], 0.7, 5.05, 7.5, 1.8, size=18, color=WHITE, dot_color=FRNDLY_YELLOW)

# Fix options
rect(s, 8.6, 2.0, 4.5, 4.8, fill=RGBColor(0x00, 0x40, 0x20))
txt(s, "The Fix (3 options)", 8.7, 2.1, 4.3, 0.5,
    size=20, bold=True, color=FRNDLY_GREEN)
bullet_box(s, [
    "① Whitelist our test account\n   (Frndly TV backend team — easiest)",
    "② Staging environment with\n   DRM disabled for QA",
    "③ Real-device testing\n   (BrowserStack — costs money)",
], 8.7, 2.7, 4.2, 4.0, size=16, color=WHITE, dot_color=FRNDLY_YELLOW, spacing=8)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Current Test Results (scorecard)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_GREEN)

txt(s, "Where We Stand Today", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

# Scorecard boxes
cards = [
    ("7", "Rows\nPASSING", FRNDLY_GREEN, FRNDLY_DARK),
    ("9", "Rows\nSKIPPED\n(DRM — fixable)", FRNDLY_YELLOW, FRNDLY_DARK),
    ("17", "Total Row\nTests", FRNDLY_BLUE, WHITE),
    ("<3s", "Avg Live TV\nLoad Time", FRNDLY_GREEN, FRNDLY_DARK),
]
for i, (number, label, bg, fg) in enumerate(cards):
    x = 0.4 + i * 3.2
    rect(s, x, 2.3, 3.0, 3.2, fill=bg)
    txt(s, number, x, 2.5, 3.0, 1.2,
        size=52, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(s, label, x, 3.8, 3.0, 1.5,
        size=17, bold=True, color=fg, align=PP_ALIGN.CENTER)

txt(s, "Once DRM is resolved — all 17 rows become green.",
    0.5, 5.9, 12, 0.6, size=22, bold=True,
    color=FRNDLY_YELLOW, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — What's Next
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
accent_bar(s, FRNDLY_GREEN)

txt(s, "What's Next", 0.5, 1.1, 12, 0.8,
    size=36, bold=True, color=WHITE)

next_steps = [
    ("✅  Done",        FRNDLY_GREEN,  "Smoke & regression tests running on GitHub Actions"),
    ("✅  Done",        FRNDLY_GREEN,  "Video recordings, TTFF measurement, HTML report"),
    ("✅  Done",        FRNDLY_GREEN,  "Live TV fully tested — all channels"),
    ("🔜  In Progress", FRNDLY_YELLOW, "Get DRM whitelist from Frndly TV backend team"),
    ("🔜  Planned",     FRNDLY_YELLOW, "Expand to mobile viewports (tablet / phone)"),
    ("🔜  Planned",     FRNDLY_YELLOW, "Nightly scheduled runs + email alerts on failure"),
    ("🔜  Planned",     FRNDLY_YELLOW, "Performance trending dashboard — TTFF over time"),
]
for i, (status, color, desc) in enumerate(next_steps):
    y = 2.1 + i * 0.72
    txt(s, status, 0.5, y, 2.5, 0.65, size=17, bold=True, color=color)
    txt(s, desc,   3.2, y, 9.8, 0.65, size=20, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Closing / Q&A
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=FRNDLY_DARK)
rect(s, 0, 0, 13.33, 0.12, fill=FRNDLY_BLUE)
rect(s, 0, 7.38, 13.33, 0.12, fill=FRNDLY_BLUE)
rect(s, 0.4, 3.4, 12.5, 0.08, fill=FRNDLY_GREEN)

txt(s, "Our subscribers deserve a product that\nworks every time they press play.",
    0.5, 1.3, 12, 1.8, size=34, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Automation is how we guarantee that.",
    0.5, 3.6, 12, 0.7, size=28, bold=True,
    color=FRNDLY_GREEN, align=PP_ALIGN.CENTER)

txt(s, "Questions?",
    0.5, 5.6, 12, 0.9, size=38, bold=True,
    color=FRNDLY_YELLOW, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/dlucero/projects/Claude-test-automation/FrndlyTV_Automation_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
